"""
SourceLens — PPTX Parser
Extracts slide content, shapes, text, and tables from PowerPoint files.
"""

from pptx import Presentation
from pptx.util import Emu
from app.models import SlideRecord, SlideShape, ShapePosition
import logging

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str) -> list[SlideRecord]:
    """
    Parse a PPTX file and extract all slide content.
    Returns a list of SlideRecord objects with shape text, tables, and positions.
    """
    prs = Presentation(file_path)
    slides = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1
        slide_id = slide.slide_id

        # Try to extract slide title
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()

        shapes = []
        all_text_parts = []

        for shape in slide.shapes:
            shape_data = _extract_shape(shape)
            if shape_data:
                shapes.append(shape_data)
                if shape_data.text:
                    all_text_parts.append(shape_data.text)

        slide_record = SlideRecord(
            slide_number=slide_number,
            slide_title=slide_title,
            slide_id=slide_id,
            shapes=shapes,
            all_text="\n\n".join(all_text_parts),
        )
        slides.append(slide_record)
        logger.info(f"Parsed slide {slide_number}: '{slide_title}' — {len(shapes)} shapes")

    return slides


def _extract_shape(shape) -> SlideShape | None:
    """Extract text and/or table data from a single shape."""
    text = ""
    table_data = None

    # Extract text from text frame
    if shape.has_text_frame:
        parts = []
        for paragraph in shape.text_frame.paragraphs:
            para_text = ""
            for run in paragraph.runs:
                para_text += run.text
            if para_text.strip():
                parts.append(para_text.strip())
        text = "\n".join(parts)

    # Extract table data
    if shape.has_table:
        table_data = []
        for row in shape.table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data.append(row_data)
        # Also flatten table to text for search
        if not text:
            table_text_parts = []
            for row in table_data:
                table_text_parts.append(" | ".join(row))
            text = "\n".join(table_text_parts)

    # Skip shapes with no content
    if not text and not table_data:
        return None

    position = ShapePosition(
        left=shape.left or 0,
        top=shape.top or 0,
        width=shape.width or 0,
        height=shape.height or 0,
    )

    return SlideShape(
        shape_id=shape.shape_id,
        shape_name=shape.name or "",
        text=text,
        table=table_data,
        position=position,
    )


def get_slide_count(file_path: str) -> int:
    """Quick slide count without full parsing."""
    prs = Presentation(file_path)
    return len(prs.slides)


def get_slide_text_preview(file_path: str, slide_number: int, max_chars: int = 500) -> str:
    """Get a text preview for a specific slide."""
    prs = Presentation(file_path)
    if slide_number < 1 or slide_number > len(prs.slides):
        return ""
    slide = prs.slides[slide_number - 1]
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
    full_text = "\n".join(parts)
    return full_text[:max_chars]
